"""
parsers/pptx_parser.py — Extract content from PowerPoint files (.pptx / .ppt).

Each slide becomes one PageContent.  Shapes are iterated in Z-order to
preserve reading sequence.  Slides with few/no text shapes but with images
are flagged as diagram pages.  Speaker notes are captured in the notes field.
"""

from __future__ import annotations

import io
import subprocess
import tempfile
from pathlib import Path
from typing import List, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from app.parsers.base_parser import BaseParser, ImageBlock, PageContent, ParserError
from app.config import settings
from app.utils.logger import get_logger

logger = get_logger(__name__)

_DIAGRAM_SHAPE_TYPES = {
    MSO_SHAPE_TYPE.AUTO_SHAPE,   # connectors, flowchart shapes
    MSO_SHAPE_TYPE.LINE,
    MSO_SHAPE_TYPE.FREEFORM,
    MSO_SHAPE_TYPE.GROUP,
}

_DIAGRAM_KEYWORDS = {
    "start", "end", "yes", "no", "decision", "process", "flow",
    "step", "if", "then", "else", "loop", "repeat", "→", "->",
    "condition", "branch", "merge", "fork", "gateway",
}


def _slide_is_diagram(shapes, text: str) -> bool:
    diagram_shapes = sum(
        1 for s in shapes if s.shape_type in _DIAGRAM_SHAPE_TYPES
    )
    kw_hits = sum(1 for kw in _DIAGRAM_KEYWORDS if kw in text.lower())
    # Many auto-shapes + few words, or keywords present → diagram
    return (diagram_shapes >= 3 and len(text.strip()) < 300) or kw_hits >= settings.diagram_keyword_threshold


def _convert_ppt_to_pptx(ppt_path: Path) -> Optional[Path]:
    """Convert legacy .ppt → .pptx using LibreOffice (if available)."""
    try:
        tmp_dir = Path(tempfile.mkdtemp())
        result = subprocess.run(
            [
                "libreoffice", "--headless",
                "--convert-to", "pptx",
                "--outdir", str(tmp_dir),
                str(ppt_path),
            ],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            converted = list(tmp_dir.glob("*.pptx"))
            if converted:
                logger.info("Converted %s → %s", ppt_path.name, converted[0].name)
                return converted[0]
        logger.warning("LibreOffice conversion failed: %s", result.stderr)
    except FileNotFoundError:
        logger.warning("LibreOffice not found; cannot convert .ppt files")
    except Exception as exc:
        logger.warning("PPT conversion error: %s", exc)
    return None


class PptxParser(BaseParser):
    """Parse .pptx (and .ppt via LibreOffice conversion) files."""

    def parse(self) -> List[PageContent]:
        target = self.file_path

        # Convert .ppt to .pptx if necessary
        if self.file_path.suffix.lower() == ".ppt":
            converted = _convert_ppt_to_pptx(self.file_path)
            if converted is None:
                raise ParserError(
                    f"Cannot parse legacy .ppt without LibreOffice: {self.file_name}"
                )
            target = converted

        try:
            prs = Presentation(str(target))
        except Exception as exc:
            raise ParserError(f"Cannot open presentation {self.file_name}: {exc}") from exc

        pages: List[PageContent] = []

        for slide_idx, slide in enumerate(prs.slides):
            page_number = slide_idx + 1
            try:
                text_parts: List[str] = []
                image_blocks: List[ImageBlock] = []
                section_title = ""

                for shape in slide.shapes:
                    # ── Text ──────────────────────────────────────────────────
                    if shape.has_text_frame:
                        shape_text = shape.text_frame.text.strip()
                        if shape_text:
                            # Heuristic: title placeholder or largest font
                            if (
                                not section_title
                                and hasattr(shape, "placeholder_format")
                                and shape.placeholder_format is not None
                                and shape.placeholder_format.idx == 0  # title
                            ):
                                section_title = shape_text[:120]
                            text_parts.append(shape_text)

                    # ── Embedded images ───────────────────────────────────────
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_blob = shape.image.blob
                            ext = shape.image.ext  # "png", "jpeg", etc.
                            mime = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                            image_blocks.append(
                                ImageBlock(
                                    image_bytes=img_blob,
                                    mime_type=mime,
                                    page_number=page_number,
                                )
                            )
                        except Exception as img_exc:
                            logger.warning(
                                "Slide %d image extraction error in %s: %s",
                                page_number, self.file_name, img_exc,
                            )

                full_text = "\n".join(text_parts)

                # ── Speaker notes ─────────────────────────────────────────────
                notes_text = ""
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    if notes_tf:
                        notes_text = notes_tf.text.strip()

                is_diagram = _slide_is_diagram(list(slide.shapes), full_text)

                pages.append(
                    PageContent(
                        page_number=page_number,
                        text=full_text,
                        images=image_blocks,
                        notes=notes_text,
                        section_title=section_title or f"Slide {page_number}",
                        is_diagram_page=is_diagram,
                        raw_metadata={
                            "source": "pptx",
                            "shape_count": len(list(slide.shapes)),
                        },
                    )
                )

            except Exception as slide_exc:
                logger.error(
                    "Error processing slide %d of %s: %s",
                    page_number, self.file_name, slide_exc,
                )
                pages.append(
                    PageContent(
                        page_number=page_number,
                        text=f"[ERROR extracting slide {page_number}: {slide_exc}]",
                        raw_metadata={"error": str(slide_exc)},
                    )
                )

        logger.info("PPTX parsed: %s — %d slide(s)", self.file_name, len(pages))
        return pages